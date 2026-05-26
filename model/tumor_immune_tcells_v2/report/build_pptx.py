"""Generate PowerPoint from report_sac_tme_state_spaces.md."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

BASE = Path(__file__).resolve().parent
IMG_DIR = BASE.parent / "img"
FIG_DIR = BASE / "figures"
OUT = BASE / "report_sac_tme_state_spaces.pptx"

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

PRIMARY = RGBColor(0x1F, 0x3A, 0x68)      # deep blue
ACCENT = RGBColor(0xC0, 0x39, 0x2B)       # deep red
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
TEXT = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, bullets, size=14, color=TEXT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = "•  " + b
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), PRIMARY)
    add_text(slide, Inches(0.4), Inches(0.12), Inches(12), Inches(0.45),
             title, size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.5), Inches(12), Inches(0.35),
                 subtitle, size=12, color=RGBColor(0xCF, 0xD8, 0xE3))
    add_rect(slide, 0, Inches(0.85), SLIDE_W, Inches(0.04), ACCENT)


def footer(slide, page, total):
    add_text(slide, Inches(0.4), Inches(7.15), Inches(8), Inches(0.3),
             "SAC-TME State Space Comparison  •  PhysiCell TME with T-cells",
             size=10, color=MUTED)
    add_text(slide, Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
             f"{page} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def fit_image(slide, path, x, y, w_max, h_max):
    """Add image fit inside the (w_max, h_max) box, centered."""
    with Image.open(path) as im:
        iw, ih = im.size
    box_ratio = w_max / h_max
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        w = w_max
        h = int(w_max / img_ratio)
    else:
        h = h_max
        w = int(h_max * img_ratio)
    px = x + (w_max - w) // 2
    py = y + (h_max - h) // 2
    slide.shapes.add_picture(str(path), px, py, width=w, height=h)


def add_table(slide, x, y, w, h, data, header=True, font_size=11,
              col_widths=None):
    rows = len(data)
    cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    tbl = tbl_shape.table
    if col_widths:
        total = sum(col_widths)
        for j, cw in enumerate(col_widths):
            tbl.columns[j].width = int(w * cw / total)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            if header and i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY
                fc = WHITE
                bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
                fc = TEXT
                bold = False
            tf = cell.text_frame
            tf.word_wrap = True
            tf.text = ""
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 or (header and i == 0) else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(font_size)
            r.font.bold = bold
            r.font.color.rgb = fc
            r.font.name = "Calibri"
    return tbl


# ---------- SLIDES ----------
slides_built = []

# Slide 1 — Title
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, PRIMARY)
add_rect(s, 0, Inches(4.6), SLIDE_W, Inches(0.05), ACCENT)
add_text(s, Inches(0.8), Inches(2.0), Inches(12), Inches(1.2),
         "SAC-TME State Space Comparison",
         size=44, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.1), Inches(12), Inches(0.8),
         "Observation design for SAC on a PhysiCell tumour micro-environment",
         size=20, color=RGBColor(0xCF, 0xD8, 0xE3))
add_text(s, Inches(0.8), Inches(4.9), Inches(12), Inches(0.5),
         "Project: SAC_ASYNC_TME_Tcells (W&B)", size=14, color=WHITE)
add_text(s, Inches(0.8), Inches(5.3), Inches(12), Inches(0.5),
         "Algorithm: Soft Actor-Critic with asynchronous parallel environments",
         size=14, color=WHITE)
add_text(s, Inches(0.8), Inches(5.7), Inches(12), Inches(0.5),
         "Environment: PhysiGym — PhysiCell TME with T-cells", size=14, color=WHITE)
add_text(s, Inches(0.8), Inches(6.5), Inches(12), Inches(0.4),
         "2026-05-18", size=14, bold=True, color=ACCENT)
slides_built.append(s)

# Slide 2 — Agenda
s = prs.slides.add_slide(BLANK)
header_bar(s, "Agenda")
items = [
    "Environment & task — PhysiCell TME with T-cells",
    "Biological model — cell types, substrates, rules",
    "Reward function & training setup",
    "Seven observation modes (S1–S5, I1–I2)",
    "Aggregated results across seeds",
    "Episode case studies (run_000143, run_000147)",
    "W&B training & test curves",
    "Key findings & recommendation",
]
add_bullets(s, Inches(1.0), Inches(1.4), Inches(11), Inches(5.5), items, size=20)
footer(s, 2, 0)
slides_built.append(s)

# Slide 3 — Environment & Task
s = prs.slides.add_slide(BLANK)
header_bar(s, "1. Environment & Task",
           "PhysiCell tumour micro-environment, 2D, 63×63 µm")
add_text(s, Inches(0.4), Inches(1.05), Inches(6.4), Inches(0.4),
         "Cell types modelled", size=16, bold=True, color=PRIMARY)
add_table(s, Inches(0.4), Inches(1.5), Inches(6.4), Inches(1.8),
          [["Cell type", "Role"],
           ["tumor", "Target — to be eliminated"],
           ["macrophage", "Immune effector (M1 / M2)"],
           ["t_cell", "Adaptive immune effector"]],
          col_widths=[2, 4], font_size=12)

add_text(s, Inches(0.4), Inches(3.5), Inches(6.4), Inches(0.4),
         "Domain & time", size=16, bold=True, color=PRIMARY)
add_table(s, Inches(0.4), Inches(3.95), Inches(6.4), Inches(2.4),
          [["Parameter", "Value"],
           ["Domain", "63 × 63 µm (2D)"],
           ["Diffusion dt", "0.01 min"],
           ["Mechanics dt", "0.1 min"],
           ["Phenotype dt", "6 min"],
           ["Max sim time", "10 080 min (7 days)"],
           ["RL gym step", "15 min"]],
          col_widths=[3, 4], font_size=11)

add_text(s, Inches(7.0), Inches(1.05), Inches(6.0), Inches(0.4),
         "Five microenvironment substrates", size=16, bold=True, color=PRIMARY)
add_table(s, Inches(7.0), Inches(1.5), Inches(6.0), Inches(4.8),
          [["Substrate", "Diff.", "Decay", "Role"],
           ["anti_tumoral_factor", "3.0", "0.001", "Attracts T-cells"],
           ["pro_tumoral_factor", "3.0", "0.001", "Promotes tumor growth"],
           ["drug_1", "0", "0.05", "Agent-controlled drug"],
           ["tumor_molecule", "2.0", "0.1", "Danger signal (tumor)"],
           ["cytokine", "0.003", "0.01", "From T-cells, kills tumor"]],
          col_widths=[3, 1, 1, 3], font_size=11)

add_text(s, Inches(7.0), Inches(6.5), Inches(6.0), Inches(0.4),
         "Action: spatially-targeted drug_1 delivery",
         size=12, bold=True, color=ACCENT)
footer(s, 3, 0)
slides_built.append(s)

# Slide 4 — Biological model: cell types
s = prs.slides.add_slide(BLANK)
header_bar(s, "2. Biological Model — Cell Properties")
add_table(s, Inches(0.4), Inches(1.1), Inches(12.5), Inches(5.6),
          [["Property", "Tumor", "T-cell", "Macrophage"],
           ["Proliferation rate", "3×10⁻⁴ / min ≈ 0.432 / day", "0 (no division)", "0"],
           ["Apoptosis rate", "1×10⁻⁶ / min (rises with cytokine)", "0 (immortal in sim)", "0"],
           ["Motility", "Immobile (speed = 0)", "Active (speed 0.1 µm/min)", "Active (speed 0.1 µm/min)"],
           ["Chemotaxis", "—", "+0.3 toward anti_tumoral_factor", "+1 toward tumor_molecule"],
           ["Secretion", "tumor_molecule @ rate 1.0", "cytokine @ rate 1.0", "pro / anti tumoral factor (rule-driven)"],
           ["Uptake", "—", "anti & pro tumoral factor (rate 1.0)", "—"],
           ["Adhesion", "All cells (affinity 1.0)", "Tumor only (1.0)", "Tumor only (1.0)"]],
          col_widths=[3, 4, 4, 4], font_size=11)
footer(s, 4, 0)
slides_built.append(s)

# Slide 5 — Cell rules (CBHG)
s = prs.slides.add_slide(BLANK)
header_bar(s, "2.4 Cell Rules (CBHG v3.0)",
           "Signal → response, applied every phenotype step (6 min)")
add_table(s, Inches(0.4), Inches(1.1), Inches(12.5), Inches(3.2),
          [["Cell", "Signal", "Effect", "Target", "Max", "Half", "n"],
           ["macrophage", "tumor_molecule ↑", "increases", "pro_tumoral_factor secretion", "10", "0.5", "1"],
           ["macrophage", "drug_1 ↑", "decreases", "pro_tumoral_factor secretion", "0", "0.5", "4"],
           ["macrophage", "drug_1 ↑", "increases", "anti_tumoral_factor secretion", "10", "0.5", "4"],
           ["macrophage", "tumor_molecule ↑", "decreases", "anti_tumoral_factor secretion", "0", "0.5", "1"],
           ["macrophage", "tumor_molecule ↑", "decreases", "migration speed", "0.0001", "0.5", "4"],
           ["tumor", "cytokine ↑", "increases", "apoptosis rate", "1", "0.5", "4"]],
          col_widths=[2, 3, 2, 4, 1, 1, 1], font_size=10)

add_text(s, Inches(0.4), Inches(4.5), Inches(12.5), Inches(0.4),
         "Biological interpretation", size=16, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.6), Inches(4.95), Inches(12.5), Inches(2.2), [
    "Tumor burden → M2-like macrophages: ↑ pro_tumoral_factor, ↓ anti_tumoral_factor, slower migration.",
    "drug_1 re-polarises macrophages to M1-like: suppresses pro_tumoral, boosts anti_tumoral — the therapeutic mechanism.",
    "T-cell cytotoxicity is delivered via cytokine: T-cells must physically reach the tumor to kill it (Hill function, n=4).",
], size=13)
footer(s, 5, 0)
slides_built.append(s)

# Slide 6 — Signalling network
s = prs.slides.add_slide(BLANK)
header_bar(s, "2.5 Complete Signalling Network")
network = (
"Tumor\n"
"  ├── secretes ─▶ tumor_molecule ──▶ macrophage chemotaxis (attracts)\n"
"  │                                ├─▶ macrophage rule: ↑ pro_tumoral_factor\n"
"  │                                ├─▶ macrophage rule: ↓ anti_tumoral_factor\n"
"  │                                └─▶ macrophage rule: ↓ migration speed\n"
"  └── responds to ◀── cytokine (T-cell) ──▶ ↑ apoptosis rate\n"
"\n"
"T-cell\n"
"  ├── secretes ─▶ cytokine ──▶ tumor apoptosis (rule)\n"
"  ├── uptakes ─── anti_tumoral_factor   (depletes pool)\n"
"  ├── uptakes ─── pro_tumoral_factor    (depletes pool)\n"
"  └── chemotaxis ─▶ toward anti_tumoral_factor (+0.3)\n"
"        ◀── indirectly REPELLED from pro_tumoral_factor zones\n"
"\n"
"Macrophage\n"
"  ├── secretes ─▶ pro_tumoral_factor  (tumor_molecule high → M2)\n"
"  ├── secretes ─▶ anti_tumoral_factor (drug_1 high → M1)\n"
"  └── chemotaxis ─▶ toward tumor_molecule\n"
"\n"
"drug_1 (agent-controlled, non-diffusing, half-life ≈ 14 min)\n"
"  └── re-polarises macrophages: ↓ pro_tumoral, ↑ anti_tumoral"
)
tb = add_text(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.9),
              network, size=12, color=TEXT, font="Consolas")
footer(s, 6, 0)
slides_built.append(s)

# Slide 7 — Why substrate observations matter
s = prs.slides.add_slide(BLANK)
header_bar(s, "2.6 Why Substrate Observations Are Crucial")
add_bullets(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5), [
    "The drug effect is indirect — the agent cannot observe tumor death per step; it must infer macrophage polarisation from pro/anti tumoral factor gradients.",
    "T-cell positioning is functionally coupled to substrate gradients — observing them lets the agent predict where cytokine-mediated killing will occur.",
    "pro_tumoral_factor is the key immunosuppression indicator — high values mean both M2 polarisation and T-cell repulsion.",
    "A scalar max(pro_tumoral_factor) (S2) loses spatial structure — the agent cannot tell WHERE the immunosuppressive zone is relative to tumor or T-cells.",
    "The 64×64 substrate channels preserve this spatial coupling — explaining why image modes generalise far better to novel layouts (circular, rectangle).",
], size=16)
footer(s, 7, 0)
slides_built.append(s)

# Slide 8 — Reward function
s = prs.slides.add_slide(BLANK)
header_bar(s, "3. Reward Function",
           "Normalised tumor-reduction signal, trajectory-sensitive")
add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
         "r(t)  =  (c_{t-1} − c_t) / expected_growth",
         size=22, bold=True, color=ACCENT, font="Consolas")
add_text(s, Inches(0.5), Inches(1.95), Inches(12), Inches(0.4),
         "where  expected_growth = c_{t-1} · (exp(λ·Δt) − 1)",
         size=14, color=TEXT, font="Consolas")
add_bullets(s, Inches(0.5), Inches(2.6), Inches(12), Inches(2.3), [
    "c_t = alive tumor count at step t",
    "λ = intrinsic tumor growth rate (PhysiCell XML, 2.99×10⁻⁴ /min)",
    "Δt = gym step duration (15 min)",
    "Reward weight: w_cell = 0.3",
], size=14)
add_text(s, Inches(0.5), Inches(4.7), Inches(12), Inches(0.4),
         "Interpretation", size=16, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.5), Inches(5.05), Inches(12), Inches(1.6), [
    "+1.0 → agent suppressed exactly the natural growth.",
    "0     → no net change.   |   negative → tumor grew.",
    "Episode terminates: c_t ≤ 3 (eradication) or c_t > 256 (uncontrolled).",
    "Metric: charts/train_return_mean50 — EMA over last 50 episodes.",
], size=14)
footer(s, 8, 0)
slides_built.append(s)

# Slide 9 — Seven state spaces
s = prs.slides.add_slide(BLANK)
header_bar(s, "4. State Spaces Compared",
           "Seven observation modes — five scalar, two image")
add_table(s, Inches(0.3), Inches(1.05), Inches(12.7), Inches(4.6),
          [["ID", "Mode", "Type", "Shape", "Description"],
           ["S1", "scalars_cells", "scalar", "(3,)", "Normalised alive count per cell type"],
           ["S2", "scalars_cells_substrates", "scalar", "(8,)", "Cell counts + max substrate concentration"],
           ["S3", "spatial_scalars_cells", "scalar", "(21,)", "Cell counts + per-type spatial statistics"],
           ["S4", "spatial_scalars_cells_substrates", "scalar", "(27,)", "Cell counts + substrate scalars + cell spatial stats"],
           ["S5", "spatial_scalars_cells_spatial_substrates", "scalar", "(39,)", "Cells + substrate scalars + substrate spatial features"],
           ["I1", "img_mc_cells", "image", "(3, 64, 64)", "Per cell type spatial density grid"],
           ["I2", "img_mc_cells_substrates", "image", "(8, 64, 64)", "Cell density grids + all 5 substrate concentration grids"]],
          col_widths=[1, 5, 1.5, 2, 6], font_size=11)
add_text(s, Inches(0.3), Inches(5.85), Inches(12.7), Inches(0.4),
         "Scalar modes use MLP policy; image modes use IMPALA-style CNN.",
         size=13, bold=True, color=MUTED)
footer(s, 9, 0)
slides_built.append(s)

# Slide 10 — Training setup
s = prs.slides.add_slide(BLANK)
header_bar(s, "5. Training Setup")
add_table(s, Inches(0.4), Inches(1.1), Inches(6.0), Inches(5.7),
          [["Parameter", "Value"],
           ["Algorithm", "SAC (Soft Actor-Critic)"],
           ["Total timesteps", "500 000"],
           ["Num parallel envs", "9"],
           ["Replay buffer size", "300 000"],
           ["Batch size", "576"],
           ["Policy LR / Q LR", "3e-4"],
           ["Gamma", "0.99"],
           ["Tau (soft update)", "0.005"],
           ["Entropy tuning", "Automatic"],
           ["Learning starts", "5 000 steps"],
           ["Test frequency", "every 4 episodes"]],
          col_widths=[3, 4], font_size=12)

add_text(s, Inches(6.7), Inches(1.05), Inches(6.4), Inches(0.4),
         "Image-mode CNN (IMPALA-style)", size=16, bold=True, color=PRIMARY)
cnn = ("Conv(in→32, 8×8, stride 4) → Mish\n"
       "Conv(32→64, 4×4, stride 2) → Mish\n"
       "Conv(64→64, 3×3, stride 1) → Mish\n"
       "Flatten → FC(256) → FC(256)")
add_text(s, Inches(6.7), Inches(1.55), Inches(6.4), Inches(1.8),
         cnn, size=13, color=TEXT, font="Consolas")

add_text(s, Inches(6.7), Inches(3.6), Inches(6.4), Inches(0.4),
         "Layouts", size=16, bold=True, color=PRIMARY)
add_bullets(s, Inches(6.7), Inches(4.0), Inches(6.4), Inches(2.5), [
    "Training: network-field generative process (Gaussian random field, ℓ = 45 µm, threshold 0.55)",
    "Test (held-out): circular, rectangle",
    "Default cells.csv: ~128 tumor, 32 T-cell, 16 macrophage",
], size=13)
footer(s, 10, 0)
slides_built.append(s)

# Slide 11 — Spatial layouts figure
s = prs.slides.add_slide(BLANK)
header_bar(s, "9.1 Spatial Layout Examples",
           "Training layouts are network-field; test layouts are held out")
fit_image(s, FIG_DIR / "fig_layouts.png",
          Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.4))
add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
         "Test layouts (rectangle, circular) are out-of-distribution: good generalisation requires policies that transfer without memorising training positions.",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)
footer(s, 11, 0)
slides_built.append(s)

# Slide 12 — Aggregated results
s = prs.slides.add_slide(BLANK)
header_bar(s, "6. Aggregated Results — by Observation Mode",
           "Only complete runs (≥ 400 k steps); multiple seeds averaged")
add_table(s, Inches(0.3), Inches(1.1), Inches(12.7), Inches(4.3),
          [["Mode", "n", "Train mean50", "Test mean50", "Test Rect", "Test Circ"],
           ["I2 — img_mc_cells_substrates", "5", "+77.2", "+56.5", "+44.5", "+97.9"],
           ["I1 — img_mc_cells", "4", "+50.3", "+20.8", "−0.6", "+5.0"],
           ["S3 — spatial_scalars_cells", "3", "+51.2", "+12.4", "+112.7", "+16.3"],
           ["S4 — spatial_scalars_cells_substrates", "3", "+29.2", "+3.4", "−22.1", "−1.2"],
           ["S5 — spatial_scalars_cells_spatial_substrates", "2", "+3.2", "+0.5", "+85.1", "−37.6"],
           ["S1 — scalars_cells", "4", "+5.9", "+17.6", "−7.7", "−20.5"],
           ["S2 — scalars_cells_substrates", "4", "−16.8", "+12.3", "+40.2", "−47.7"]],
          col_widths=[5, 1, 2, 2, 2, 2], font_size=11)
add_text(s, Inches(0.3), Inches(5.6), Inches(12.7), Inches(1.4),
         "I2 (image + substrates) leads on every metric except Test Rect, where S3 is helped by a single lucky seed (#64).\n"
         "I2 also has the lowest seed-to-seed variance — both highest mean and most reliable.",
         size=14, bold=True, color=PRIMARY)
footer(s, 12, 0)
slides_built.append(s)

# Slide 13 — Per-seed I2 + I1
s = prs.slides.add_slide(BLANK)
header_bar(s, "6. Per-Seed Breakdown — Image Modes")
add_text(s, Inches(0.3), Inches(1.0), Inches(6.3), Inches(0.4),
         "I2 — img_mc_cells_substrates (6ch)", size=14, bold=True, color=ACCENT)
add_table(s, Inches(0.3), Inches(1.45), Inches(6.3), Inches(2.5),
          [["Seed", "Steps", "Train50", "Test50", "Rect", "Circ"],
           ["1",   "499 356", "+78.0", "+58.2", "+4.4",   "+106.7"],
           ["32",  "529 501", "+79.7", "+69.9", "+4.9",   "+61.4"],
           ["64",  "410 875", "+86.9", "+47.7", "+109.6", "+33.7"],
           ["64",  "499 977", "+64.3", "+42.3", "+47.3",  "+128.9"],
           ["128", "499 914", "+77.1", "+64.5", "+56.1",  "+158.8"]],
          col_widths=[1, 2, 1.5, 1.5, 1.5, 1.5], font_size=10)

add_text(s, Inches(6.8), Inches(1.0), Inches(6.3), Inches(0.4),
         "I1 — img_mc_cells (3ch)", size=14, bold=True, color=ACCENT)
add_table(s, Inches(6.8), Inches(1.45), Inches(6.3), Inches(2.5),
          [["Seed", "Steps", "Train50", "Test50", "Rect", "Circ"],
           ["1",   "499 884", "+39.8", "+22.4", "−28.7", "−28.1"],
           ["32",  "526 602", "+42.4", "+27.8", "−17.4", "−43.9"],
           ["64",  "499 769", "+51.3", "+14.5", "−0.5",  "−2.8"],
           ["128", "499 632", "+67.5", "+18.3", "+44.2", "+94.8"]],
          col_widths=[1, 2, 1.5, 1.5, 1.5, 1.5], font_size=10)

add_text(s, Inches(0.3), Inches(4.3), Inches(12.8), Inches(0.4),
         "S3 — spatial_scalars_cells (21d)", size=14, bold=True, color=ACCENT)
add_table(s, Inches(0.3), Inches(4.75), Inches(6.3), Inches(1.8),
          [["Seed", "Steps", "Train50", "Test50", "Rect", "Circ"],
           ["1",   "499 976", "+16.7", "−28.9", "+32.2",  "−81.3"],
           ["64",  "499 996", "+59.8", "+27.3", "+223.1", "+91.1"],
           ["128", "499 990", "+77.1", "+38.8", "+82.9",  "+39.2"]],
          col_widths=[1, 2, 1.5, 1.5, 1.5, 1.5], font_size=10)

add_text(s, Inches(6.8), Inches(4.3), Inches(6.3), Inches(0.4),
         "Key takeaways", size=14, bold=True, color=PRIMARY)
add_bullets(s, Inches(6.8), Inches(4.75), Inches(6.3), Inches(2.3), [
    "I2 train50 in tight band [+64, +87] — most consistent mode.",
    "I1 collapses on rect/circ for 3 of 4 seeds (overfits training layout).",
    "S3 seed #64 (+223 rect) is an outlier, not a robust trend.",
], size=12)
footer(s, 13, 0)
slides_built.append(s)

# Slide 14 — Per-seed scalar modes S1/S2/S4/S5
s = prs.slides.add_slide(BLANK)
header_bar(s, "6. Per-Seed Breakdown — Scalar Modes")
add_text(s, Inches(0.3), Inches(1.0), Inches(6.3), Inches(0.4),
         "S4 — spatial_scalars_cells_substrates (27d)",
         size=13, bold=True, color=ACCENT)
add_table(s, Inches(0.3), Inches(1.45), Inches(6.3), Inches(1.8),
          [["Seed", "Steps", "Train50", "Test50", "Rect", "Circ"],
           ["1",   "499 981", "−1.7",  "−7.2",  "−107.0", "−57.7"],
           ["64",  "499 951", "−2.3",  "−28.5", "−5.4",   "−15.5"],
           ["128", "499 938", "+91.5", "+46.0", "+46.0",  "+69.5"]],
          col_widths=[1, 2, 1.5, 1.5, 1.5, 1.5], font_size=10)

add_text(s, Inches(6.8), Inches(1.0), Inches(6.3), Inches(0.4),
         "S5 — spatial_scalars_cells_spatial_substrates (39d)",
         size=13, bold=True, color=ACCENT)
add_table(s, Inches(6.8), Inches(1.45), Inches(6.3), Inches(1.8),
          [["Seed", "Steps", "Train50", "Test50", "Rect", "Circ"],
           ["1",   "499 979", "+10.4", "+1.2",  "+6.0",   "−67.5"],
           ["64",  "499 979", "−4.1",  "−0.1",  "+164.2", "−7.6"]],
          col_widths=[1, 2, 1.5, 1.5, 1.5, 1.5], font_size=10)

add_text(s, Inches(0.3), Inches(3.55), Inches(6.3), Inches(0.4),
         "S1 — scalars_cells (3d)", size=13, bold=True, color=ACCENT)
add_table(s, Inches(0.3), Inches(4.0), Inches(6.3), Inches(2.2),
          [["Seed", "Steps", "Train50", "Test50", "Rect", "Circ"],
           ["1",   "499 925", "+37.2", "+24.0", "−1.5",  "−21.8"],
           ["32",  "527 895", "+6.2",  "+23.3", "+30.3", "+46.7"],
           ["64",  "499 988", "−20.3", "+6.9",  "−52.6", "−33.1"],
           ["128", "499 938", "+0.4",  "+16.3", "−6.9",  "−74.0"]],
          col_widths=[1, 2, 1.5, 1.5, 1.5, 1.5], font_size=10)

add_text(s, Inches(6.8), Inches(3.55), Inches(6.3), Inches(0.4),
         "S2 — scalars_cells_substrates (8d)", size=13, bold=True, color=ACCENT)
add_table(s, Inches(6.8), Inches(4.0), Inches(6.3), Inches(2.2),
          [["Seed", "Steps", "Train50", "Test50", "Rect", "Circ"],
           ["1",   "499 939", "+1.7",  "+26.8", "+144.6", "−0.6"],
           ["32",  "528 409", "−79.2", "−24.6", "−4.8",   "−90.8"],
           ["64",  "499 942", "+25.3", "+18.8", "+8.0",   "−44.4"],
           ["128", "499 959", "−15.1", "+28.1", "+12.9",  "−54.8"]],
          col_widths=[1, 2, 1.5, 1.5, 1.5, 1.5], font_size=10)
footer(s, 14, 0)
slides_built.append(s)

# Slide 15 — Finding 7.1 Image > scalar
s = prs.slides.add_slide(BLANK)
header_bar(s, "7.1 Image-Based Modes Outperform Scalar Modes")
add_bullets(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5), [
    "I2 (img_mc_cells_substrates) is the best across every metric: train +77.2, test +56.5 (>2× the next mode).",
    "Low seed variance: all I2 seeds land in [+64, +87] vs S4 ranging from −2.3 to +91.5.",
    "I1 (img_mc_cells) performs well on training (+50.3) but collapses on held-out test layouts — no substrate channels means no info to generalise.",
    "Adding substrate channels (I1 → I2) gains +27 train, +36 test — substrate visibility closes the credit-assignment loop.",
    "All scalar modes (S1–S5) plateau below +35 on training. None matches even I1's training mean.",
], size=16)
footer(s, 15, 0)
slides_built.append(s)

# Slide 16 — Finding 7.3 + 7.4 — more is not better, seed sensitivity
s = prs.slides.add_slide(BLANK)
header_bar(s, "7.3–7.4  More Features ≠ Better  •  Seed Sensitivity")
add_text(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.4),
         "Scalar complexity vs train return (counter-intuitive)",
         size=15, bold=True, color=PRIMARY)
add_table(s, Inches(0.4), Inches(1.5), Inches(12.5), Inches(1.2),
          [["S1 (3d)", "S2 (8d)", "S3 (21d)", "S4 (27d)", "S5 (39d)"],
           ["+5.9", "−16.8", "+51.2", "+29.2", "+3.2"]],
          col_widths=[1, 1, 1, 1, 1], font_size=14)

add_text(s, Inches(0.4), Inches(3.0), Inches(12.5), Inches(0.4),
         "Seed-to-seed standard deviation of train_return_mean50",
         size=15, bold=True, color=PRIMARY)
add_table(s, Inches(0.4), Inches(3.45), Inches(8.0), Inches(2.4),
          [["Mode", "Std across seeds"],
           ["img_mc_cells_substrates  (I2)", "~8.5"],
           ["img_mc_cells  (I1)", "~11.3"],
           ["spatial_scalars_cells  (S3)", "~25.8"],
           ["scalars_cells  (S1)", "~22.0"],
           ["scalars_cells_substrates  (S2)", "~38.3"]],
          col_widths=[4, 2], font_size=12)

add_bullets(s, Inches(8.8), Inches(3.45), Inches(4.3), Inches(3.4), [
    "Adding substrate scalars hurts: noise dominates the signal.",
    "S5 (39d) underperforms S3 (21d) — the MLP cannot exploit the bigger scalar vector in 500k steps.",
    "Image modes are 3–4× more stable across seeds — CNN inductive bias matches the physical TME structure.",
], size=12)
footer(s, 16, 0)
slides_built.append(s)

# Slide 17 — Episode run_000143 figure
s = prs.slides.add_slide(BLANK)
header_bar(s, "10. Episode Comparison — run_000143",
           "Same matched episode across I2, I1, S3 (best three modes)")
fit_image(s, FIG_DIR / "fig_episode_comparison.png",
          Inches(0.4), Inches(1.0), Inches(8.5), Inches(5.8))
add_text(s, Inches(9.1), Inches(1.0), Inches(4.0), Inches(0.4),
         "Episode statistics", size=14, bold=True, color=PRIMARY)
add_table(s, Inches(9.1), Inches(1.45), Inches(4.0), Inches(4.8),
          [["Metric", "I2", "I1", "S3"],
           ["Final tumor", "4", "112", "14"],
           ["Cum. return", "+175.1", "−7.4", "+106.3"],
           ["Cum. dose", "35.42", "12.16", "31.05"],
           ["Peak dose/step", "0.813", "0.451", "0.484"],
           ["Length", "480", "480", "480"]],
          col_widths=[3, 2, 2, 2], font_size=10)
add_text(s, Inches(9.1), Inches(6.3), Inches(4.0), Inches(0.7),
         "I2 reaches near-eradication;\nI1 under-medicates;\nS3 stalls before final push.",
         size=11, bold=True, color=ACCENT)
footer(s, 17, 0)
slides_built.append(s)

# Slide 18 — run_000143 interpretation
s = prs.slides.add_slide(BLANK)
header_bar(s, "10.2 run_000143 — Why I2 Wins")
add_bullets(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5), [
    "Cumulative return — I2 climbs steadily to +175 (near-eradication); S3 reaches +106 then flattens at 14 tumor cells; I1 dips negative and ends with 112 cells (≈ untreated).",
    "Cumulative dose — I2 and S3 both apply ~31–35 units; I1 under-medicates at 12. Without substrate feedback, I1 can't tell if a dose worked, so it converges to a conservative low-dose policy.",
    "Dose per step — I2 shows pronounced late-episode spikes (steps 300–480), driving the final push to eradication. S3 distributes dosing more evenly without substrate-guided targeting.",
    "Substrate channels close the feedback loop: I2 observes pro/anti tumoral factor maps → knows if macrophages were re-polarised; I1 has only cell counts (too noisy for credit assignment).",
    "Spatial targeting requires spatial state. I2 sees WHERE factors concentrate → can place drug at the immunosuppressive zone. S3 has cell centroids but no substrate coupling.",
], size=14)
footer(s, 18, 0)
slides_built.append(s)

# Slide 19 — Episode run_000147 figure
s = prs.slides.add_slide(BLANK)
header_bar(s, "11. Episode Comparison — run_000147 (seed 128)",
           "A harder episode: all three agents face tumor rebound")
fit_image(s, FIG_DIR / "fig_episode_comparison_2.png",
          Inches(0.4), Inches(1.0), Inches(8.5), Inches(5.8))
add_text(s, Inches(9.1), Inches(1.0), Inches(4.0), Inches(0.4),
         "Episode statistics", size=14, bold=True, color=PRIMARY)
add_table(s, Inches(9.1), Inches(1.45), Inches(4.0), Inches(5.0),
          [["Metric", "I2", "I1", "S3"],
           ["Final tumor", "31", "64", "43"],
           ["Cum. return", "+53.4", "+35.0", "+32.3"],
           ["Peak return", "+63.1", "+58.2", "+66.6"],
           ["Cum. dose", "35.69", "7.16", "34.92"],
           ["Peak dose/step", "0.834", "0.563", "0.496"]],
          col_widths=[3, 2, 2, 2], font_size=10)
add_text(s, Inches(9.1), Inches(6.5), Inches(4.0), Inches(0.5),
         "I2 still leads despite tumor rebound.",
         size=11, bold=True, color=ACCENT)
footer(s, 19, 0)
slides_built.append(s)

# Slide 20 — run_000147 interpretation + comparison
s = prs.slides.add_slide(BLANK)
header_bar(s, "11.2–11.3  run_000147 — Rebound Dynamics")
add_bullets(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(3.0), [
    "All three agents build positive return through step ~330, then the slope flattens or turns negative — the tumor rebounds after initial suppression.",
    "I2 maintains the lead and declines least after peak (+63.1 → +53.4). S3 peaks highest (+66.6) but loses control fastest (→ +32.3).",
    "I1 again under-medicates (7.2 vs ~35 units for I2/S3) — no substrate feedback ⇒ conservative dosing.",
    "Equal cumulative dose (I2 vs S3) but better outcome for I2 ⇒ WHERE the drug goes matters as much as HOW MUCH.",
], size=13)

add_text(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.4),
         "run_000143 vs run_000147", size=14, bold=True, color=PRIMARY)
add_table(s, Inches(0.5), Inches(4.65), Inches(12.3), Inches(2.1),
          [["", "run_000143 (seed 64)", "run_000147 (seed 128)"],
           ["I2 final tumor",      "4 (near-eradication)", "31 (rebound)"],
           ["I2 cumulative return","+175.1",                "+53.4"],
           ["I1 final tumor",      "112",                  "64"],
           ["S3 final tumor",      "14",                   "43"],
           ["Return rollback?",    "No",                   "Yes"]],
          col_widths=[3, 4, 4], font_size=11)
footer(s, 20, 0)
slides_built.append(s)

# Slide 21 — Observations visualisation (I2)
s = prs.slides.add_slide(BLANK)
header_bar(s, "12.1 What the I2 Agent Sees",
           "8-channel observation at six time steps of one episode")
xs = [Inches(0.3), Inches(4.65), Inches(8.95)]
ys_top = Inches(1.05)
ys_bot = Inches(4.15)
w_img = Inches(4.2); h_img = Inches(3.0)
for i, step in enumerate([0, 20, 40]):
    fit_image(s, IMG_DIR / f"observation_{step}.png", xs[i], ys_top, w_img, h_img)
    add_text(s, xs[i], ys_top + h_img - Inches(0.05),
             w_img, Inches(0.3), f"Step {step}",
             size=11, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
for i, step in enumerate([60, 80, 100]):
    fit_image(s, IMG_DIR / f"observation_{step}.png", xs[i], ys_bot, w_img, h_img)
    add_text(s, xs[i], ys_bot + h_img - Inches(0.05),
             w_img, Inches(0.3), f"Step {step}",
             size=11, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.3), Inches(7.05), Inches(12.7), Inches(0.4),
         "CH 0–2: tumor / T-cell / macrophage density  •  CH 3–7: anti_tumoral, pro_tumoral, drug_1, tumor_molecule, cytokine (all 5 substrates included)",
         size=10, color=MUTED, align=PP_ALIGN.CENTER)
footer(s, 21, 0)
slides_built.append(s)

# Slide 22 — what each state space sees
s = prs.slides.add_slide(BLANK)
header_bar(s, "12.1 What Each State Space Actually Sees")
add_table(s, Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.0),
          [["State space", "What the agent sees", "What is missing"],
           ["I2 — img_mc_cells_substrates (8ch, 64×64)",
            "All 3 cell types + all 5 substrate maps (anti/pro tumoral, drug_1, tumor_molecule, cytokine)",
            "Nothing — full spatial state"],
           ["I1 — img_mc_cells (3ch, 64×64)",
            "Spatial layout of tumor, T-cell, macrophage only",
            "All substrate maps — no drug/immune feedback"],
           ["S3 — spatial_scalars_cells (21-dim)",
            "Per cell type: presence, x/y mean & std, dist_to_center",
            "All substrate information; only centroid/spread of cells"],
           ["S1 — scalars_cells (3-dim)",
            "Normalised alive count per cell type",
            "Spatial structure, substrate gradients — minimal signal"]],
          col_widths=[3, 5, 4], font_size=11)
add_text(s, Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.7),
         "CH 4 (pro_tumoral_factor) shows the immunosuppressive zone — precisely where the drug must go to repolarise macrophages.\n"
         "Without this channel, the agent cannot close the dose → polarisation feedback loop.",
         size=12, bold=True, color=ACCENT)
footer(s, 22, 0)
slides_built.append(s)

# Slide 23 — Training curve
s = prs.slides.add_slide(BLANK)
header_bar(s, "12.2 Training Return — All Modes",
           "charts/train_return_mean50, smoothed")
fit_image(s, IMG_DIR / "Section-2-Panel-3-uwnnqvhkz.png",
          Inches(0.5), Inches(1.0), Inches(8.5), Inches(5.8))
add_bullets(s, Inches(9.2), Inches(1.2), Inches(4.0), Inches(5.5), [
    "I2 (red dashed) separates from all others after ~50 k steps.",
    "Reaches ~75–80 by 500 k steps.",
    "I1 (dark green) converges to ~47.",
    "All scalar modes plateau below 35.",
    "S3 (blue) is the most stable scalar.",
    "S4/S5 do not beat S3 despite higher dim.",
], size=12)
footer(s, 23, 0)
slides_built.append(s)

# Slide 24 — Training std
s = prs.slides.add_slide(BLANK)
header_bar(s, "12.2 Training Return Std",
           "Higher std = wider behavioural range")
fit_image(s, IMG_DIR / "Section-2-Panel-2-g0ddriu6t.png",
          Inches(0.5), Inches(1.0), Inches(8.5), Inches(5.8))
add_bullets(s, Inches(9.2), Inches(1.2), Inches(4.0), Inches(5.5), [
    "I2 and I1 keep std ~60–65 throughout.",
    "Reflects diverse outcomes per episode — some near-eradication, some rebound.",
    "Scalar modes decay to low std — policies converge to a narrow, less ambitious behaviour.",
], size=12)
footer(s, 24, 0)
slides_built.append(s)

# Slide 25 — Test rect + circ raw
s = prs.slides.add_slide(BLANK)
header_bar(s, "12.2 Test Return — Held-out Layouts (raw)",
           "Left: rectangle  •  Right: circular")
fit_image(s, IMG_DIR / "Section-2-Panel-0-t6c8f8adz.png",
          Inches(0.2), Inches(1.0), Inches(6.5), Inches(5.5))
fit_image(s, IMG_DIR / "Section-2-Panel-1-v5gma8p0j.png",
          Inches(6.8), Inches(1.0), Inches(6.5), Inches(5.5))
add_text(s, Inches(0.4), Inches(6.55), Inches(12.7), Inches(0.4),
         "I2 leads with the highest and most consistent test returns; I1 is volatile; scalar modes cluster near zero. Circular is harder than rectangle.",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)
footer(s, 25, 0)
slides_built.append(s)

# Slide 26 — Test return mean50 + std
s = prs.slides.add_slide(BLANK)
header_bar(s, "12.2 Test Return — Smoothed and Std",
           "Left: test_return_mean50  •  Right: test_return_std")
fit_image(s, IMG_DIR / "Section-2-Panel-6-42sh5v4hc.png",
          Inches(0.2), Inches(1.0), Inches(6.5), Inches(5.5))
fit_image(s, IMG_DIR / "Section-2-Panel-4-q55wfb7fm.png",
          Inches(6.8), Inches(1.0), Inches(6.5), Inches(5.5))
add_text(s, Inches(0.4), Inches(6.55), Inches(12.7), Inches(0.4),
         "I2 climbs to ~55–60 on test mean50. High I2 test std is capability (near-eradication bursts), not instability — scalar modes converge to low, predictable returns.",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)
footer(s, 26, 0)
slides_built.append(s)

# Slide 27 — Summary
s = prs.slides.add_slide(BLANK)
header_bar(s, "8. Summary of Findings")
add_table(s, Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.4),
          [["Finding", "Conclusion"],
           ["Best overall state space", "img_mc_cells_substrates  (6-ch image, 64×64)"],
           ["Best scalar state space", "spatial_scalars_cells  (21d) — but high variance"],
           ["Worst state spaces",      "Pure global scalars: scalars_cells, scalars_cells_substrates"],
           ["Image vs scalar gap",     "+71 pts on train_mean50 (I2 vs S1); +39 pts on test_mean50"],
           ["Substrate info impact",   "Critical: I2 > I1 by +27 train, +36 test"],
           ["Spatial features value",  "Helpful in training (S3 > S1), inconsistent on test"],
           ["More is not always better","S5 (39d) ≈ S1 (3d) due to optimisation difficulty"],
           ["Seed robustness",         "Image modes ~3–4× more consistent than scalar modes"]],
          col_widths=[4, 8], font_size=12)
footer(s, 27, 0)
slides_built.append(s)

# Slide 28 — Recommendation
s = prs.slides.add_slide(BLANK)
header_bar(s, "Recommendation")
add_rect(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.5), PRIMARY)
add_text(s, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.6),
         "Default state representation for the PhysiCell TME task:",
         size=16, color=WHITE)
add_text(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.8),
         "img_mc_cells_substrates  (I2)",
         size=32, bold=True, color=WHITE)

add_text(s, Inches(0.5), Inches(2.95), Inches(12.3), Inches(0.4),
         "Provides", size=16, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.7), Inches(3.4), Inches(12), Inches(2.6), [
    "Full spatial layout of all cell types",
    "Full spatial layout of microenvironmental substrates",
    "Consistent performance across seeds (std ~8.5)",
    "Strong generalisation to unseen layouts (circular, rectangle)",
], size=15)
add_text(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.4),
         "Fallback when compute is constrained", size=14, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(6.5), Inches(12.3), Inches(0.5),
         "spatial_scalars_cells (S3) is the best scalar alternative — but needs more seeds for reliable performance.",
         size=13, color=TEXT)
footer(s, 28, 0)
slides_built.append(s)

# Final pass: update page numbers
total = len(slides_built)
for i, sl in enumerate(slides_built):
    if i == 0:
        continue
    for shape in sl.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.endswith(" / 0") or run.text == f"{i+1} / 0":
                        run.text = f"{i+1} / {total}"

prs.save(OUT)
print(f"Wrote {OUT}  ({total} slides)")
