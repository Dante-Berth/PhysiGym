#!/usr/bin/env python3
"""Build a PowerPoint deck from the PhysiGym EWRL 2026 paper."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "img")
VID = os.path.join(HERE, "..", "..", "video")  # model/tumor_immune_tcells_v2/video

# ---- palette ----
DARK = RGBColor(0x1A, 0x23, 0x3B)      # deep navy
ACCENT = RGBColor(0xC0, 0x39, 0x2B)    # paper red (I2 colour)
ACCENT2 = RGBColor(0x2E, 0x86, 0xC1)   # blue
GREEN = RGBColor(0x27, 0xAE, 0x60)
GREY = RGBColor(0x5D, 0x6D, 0x7E)
LIGHT = RGBColor(0xF4, 0xF6, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def fill(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill_color=None, line_color=None, line_w=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill_color is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill_color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = line_w or Pt(1)
    sh.shadow.inherit = False
    return sh


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt, size, bold, color, italic)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (txt, size, bold, color, *rest) in para:
            italic = rest[0] if rest else False
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color
            r.font.name = "Calibri"
    return tb


def pic_fit(slide, path, l, t, max_w, max_h, center=True):
    im = Image.open(path)
    iw, ih = im.size
    ar = iw / ih
    boxar = max_w / max_h
    if ar > boxar:
        w = max_w
        h = int(max_w / ar)
    else:
        h = max_h
        w = int(max_h * ar)
    if center:
        l = l + (max_w - w) // 2
        t = t + (max_h - h) // 2
    return slide.shapes.add_picture(path, l, t, width=w, height=h)


def header(slide, kicker, title, color=DARK):
    box(slide, 0, 0, SW, Inches(1.15), fill_color=color)
    box(slide, Inches(0.0), Inches(1.15), SW, Pt(4), fill_color=ACCENT)
    text(slide, Inches(0.55), Inches(0.12), SW - Inches(1.1), Inches(0.95),
         [[(kicker, 12, True, ACCENT)],
          [(title, 26, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE, space_after=2)


def footer(slide, n):
    text(slide, Inches(0.4), SH - Inches(0.45), Inches(8), Inches(0.35),
         [[("PhysiGym  Bridging Agent-Based Biological Simulation and RL", 9, False, GREY)]])
    text(slide, SW - Inches(1.2), SH - Inches(0.45), Inches(0.8), Inches(0.35),
         [[(str(n), 9, False, GREY)]], align=PP_ALIGN.RIGHT)


# ============================================================ 1. TITLE
s = add_slide()
fill(s, DARK)
box(s, 0, Inches(2.55), SW, Pt(4), fill_color=ACCENT)
text(s, Inches(0.8), Inches(0.9), SW - Inches(1.6), Inches(1.6),
     [[("PhysiGym", 54, True, WHITE)],
      [("Bridging Agent-Based Biological Simulation and Reinforcement Learning", 26, True, RGBColor(0xD6,0xDB,0xE3))]],
     space_after=8)
text(s, Inches(0.8), Inches(2.75), SW - Inches(1.6), Inches(0.7),
     [[("A State-Space Study on Tumor Microenvironment Control", 18, False, ACCENT2, True)]])
text(s, Inches(0.8), Inches(3.7), SW - Inches(1.6), Inches(1.6),
     [[("Alexandre Bertin·  Elmar Bucher·  Owen Griere  ·  Marcelo Hurtado  ·  Heber Rocha", 14, False, RGBColor(0xC8,0xCF,0xDA))],
      [("Randy Heiland  ·  Paul Macklin  ·  Vincent François-Lavet  ·  Emmanuel Rachelson  ·  Vera Pancaldi", 14, False, RGBColor(0xC8,0xCF,0xDA))]],
     space_after=4)
text(s, Inches(0.8), Inches(4.9), SW - Inches(1.6), Inches(1.3),
     [[("CRCT Toulouse  ·  ISAE-Supaéro  ·  Indiana University  ·  VU Amsterdam", 12, False, GREY)],
      [("EWRL 2026", 13, True, ACCENT)]], space_after=6)
text(s, Inches(0.8), Inches(6.4), SW - Inches(1.6), Inches(0.5),
     [[("Open source:  https://anonymous.4open.science/r/PhysiGym-76A7", 12, False, ACCENT2)]])

# ============================================================ 2. MOTIVATION
s = add_slide(); fill(s, WHITE)
header(s, "MOTIVATION", "Why connect biological simulation to RL?")
bullets = [
    ("Agent-based models (ABMs) ", "simulate biological systems where cells interact dynamically by rules — capturing spatial & cellular heterogeneity that ODEs cannot."),
    ("Dynamic treatment regimes ", "(adaptive, sequential interventions) need more than simulation — they need a control framework."),
    ("Reinforcement learning ", "learns policies through interaction — a natural fit for adaptive control of biological simulations."),
    ("The missing piece: ", "a standardized interface between a high-fidelity ABM and a modern RL API."),
]
y = Inches(1.55)
for head, body in bullets:
    box(s, Inches(0.7), y + Inches(0.07), Inches(0.18), Inches(0.18), fill_color=ACCENT)
    text(s, Inches(1.05), y, Inches(11.4), Inches(0.9),
         [[(head, 17, True, DARK), (body, 17, False, RGBColor(0x33,0x3a,0x44))]],
         line_spacing=1.0)
    y += Inches(1.0)
box(s, Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.85), fill_color=LIGHT)
text(s, Inches(0.95), Inches(6.3), Inches(11.4), Inches(0.75),
     [[("We introduce ", 16, False, DARK), ("PhysiGym", 16, True, ACCENT),
       (": it bridges PhysiCell (ABM) and Gymnasium (RL), turning any PhysiCell model into an RL environment with minimal code changes.", 16, False, DARK)]],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s, 2)

# ============================================================ 3. PHYSIGYM DESIGN
s = add_slide(); fill(s, WHITE)
header(s, "FRAMEWORK", "PhysiGym: design & implementation")
# left column - PhysiCell + Gym
text(s, Inches(0.6), Inches(1.45), Inches(6.0), Inches(0.5),
     [[("PhysiCell  +  Gymnasium", 19, True, DARK)]])
items = [
    ("PhysiCell", "Open-source C++ ABM framework. Cells = agents with type-specific rules; substrates (oxygen, cytokines) diffuse via the BioFVM solver. 2D/3D, off-lattice, multiscale."),
    ("Gymnasium", "Standard Python interface for MDPs (state S, action A, transition T, reward R). Promotes modular, reusable RL code."),
]
y = Inches(2.0)
for h, b in items:
    text(s, Inches(0.6), y, Inches(6.0), Inches(1.4),
         [[(h, 16, True, ACCENT2)], [(b, 13.5, False, RGBColor(0x33,0x3a,0x44))]], space_after=2)
    y += Inches(1.5)
# right column - how it connects
box(s, Inches(6.95), Inches(1.45), Inches(5.75), Inches(5.3), fill_color=LIGHT)
text(s, Inches(7.2), Inches(1.6), Inches(5.3), Inches(0.5),
     [[("How the bridge works", 17, True, DARK)]])
conn = [
    "Splits the C++ simulation loop into 3 Python-callable hooks: physicell_start, physicell_step, physicell_stop.",
    "Three read functions return state: get_cell (positions), get_microenv (substrate fields), get_graph (neighbourhoods).",
    "Configured via PhysiCell user parameters & custom data — no recompilation; works with PhysiCell Studio.",
    "Subclass CorePhysiCellEnv to define observation space, action space & reward. The base class handles all Gym bookkeeping.",
    "Tested on Linux, WSL, macOS with CI via GitHub Actions.",
]
y = Inches(2.15)
for c in conn:
    box(s, Inches(7.2), y + Inches(0.06), Inches(0.13), Inches(0.13), fill_color=ACCENT)
    text(s, Inches(7.5), y, Inches(5.0), Inches(0.95),
         [[(c, 13, False, RGBColor(0x2b,0x32,0x3c))]], line_spacing=1.0)
    y += Inches(0.92)
footer(s, 3)

# ============================================================ 4. TME TOY MODEL + network
s = add_slide(); fill(s, WHITE)
header(s, "CASE STUDY", "A tumor microenvironment (TME) toy model")
text(s, Inches(0.6), Inches(1.4), Inches(5.2), Inches(0.6),
     [[("Small, cheap, but spatially explicit", 16, True, DARK)]])
text(s, Inches(0.6), Inches(1.95), Inches(5.25), Inches(1.2),
     [[("64×64 µm 2D domain · 3 cell types · 5 substrates. Gym step Δt = 15 min, episodes up to 10,080 min (7 days).", 13.5, False, RGBColor(0x33,0x3a,0x44))]], line_spacing=1.05)
cells = [
    ("Tumor cells (Target)", "Immobile, dividing. Secrete a signal that suppresses local immunity.", ACCENT),
    ("Macrophages (Mediators)", "Mobile. Tumor signal “hijacks” them: anti- → pro-tumoral.", RGBColor(0xE6,0x7E,0x22)),
    ("T-cells (Effectors)", "Mobile, can trigger tumor apoptosis — but repelled by pro-tumoral macrophages.", ACCENT2),
]
y = Inches(3.1)
for h, b, c in cells:
    box(s, Inches(0.6), y, Inches(0.12), Inches(0.85), fill_color=c)
    text(s, Inches(0.85), y, Inches(5.0), Inches(0.9),
         [[(h, 13.5, True, c)], [(b, 12, False, RGBColor(0x33,0x3a,0x44))]], space_after=1)
    y += Inches(1.0)
box(s, Inches(0.6), Inches(6.25), Inches(5.45), Inches(0.85), fill_color=LIGHT)
text(s, Inches(0.8), Inches(6.3), Inches(5.1), Inches(0.75),
     [[("The control challenge: ", 12.5, True, ACCENT),
       ("drug_1 has zero direct cytotoxicity — the agent must re-program macrophages to clear a path for T-cells.", 12.5, False, DARK)]],
     anchor=MSO_ANCHOR.MIDDLE)
# network figure
pic_fit(s, os.path.join(IMG, "fig_tme-1.png"), Inches(6.3), Inches(1.5), Inches(6.5), Inches(5.0))
text(s, Inches(6.3), Inches(6.55), Inches(6.5), Inches(0.4),
     [[("TME signalling network — green = activate/attract, red = suppress/repel, dashed = drug intervention.", 10, False, GREY, True)]], align=PP_ALIGN.CENTER)
footer(s, 4)

# ============================================================ 5. RL SETUP
s = add_slide(); fill(s, WHITE)
header(s, "RL FORMULATION", "Action, reward & the question we study")
# Action
box(s, Inches(0.6), Inches(1.5), Inches(5.9), Inches(2.0), fill_color=LIGHT)
text(s, Inches(0.85), Inches(1.65), Inches(5.4), Inches(1.8),
     [[("Action space", 17, True, ACCENT2)],
      [("aₜ = (dose, x, y, r)  ∈  [0,1]⁴", 15, True, DARK)],
      [("dose = drug concentration applied", 13, False, RGBColor(0x33,0x3a,0x44))],
      [("(x, y) = injection centre (normalised)", 13, False, RGBColor(0x33,0x3a,0x44))],
      [("r = injection radius (norm. by diagonal)", 13, False, RGBColor(0x33,0x3a,0x44))]], space_after=4)
# Reward
box(s, Inches(6.8), Inches(1.5), Inches(5.9), Inches(2.0), fill_color=LIGHT)
text(s, Inches(7.05), Inches(1.65), Inches(5.4), Inches(1.8),
     [[("Reward", 17, True, ACCENT2)],
      [("rₜ = w_cell · (cₜ₋₁−cₜ)/[cₜ₋₁(e^{λΔt}−1)]  −  dₜ", 14, True, DARK)],
      [("Normalised tumor-suppression signal minus a drug-cost penalty.", 13, False, RGBColor(0x33,0x3a,0x44))],
      [("w_cell = 0.3 ;  dₜ = normalised dose spent.", 13, False, RGBColor(0x33,0x3a,0x44))]], space_after=4)
# Objective / question
box(s, Inches(0.6), Inches(3.75), Inches(12.1), Inches(1.15), fill_color=DARK)
text(s, Inches(0.9), Inches(3.85), Inches(11.5), Inches(1.0),
     [[("Objective:", 15, True, ACCENT), ("  eradicate the tumor while minimizing drug exposure.", 15, False, WHITE)],
      [("The research question:", 15, True, ACCENT), ("  among the observation spaces PhysiGym can expose, which help the agent learn & generalize?", 15, False, WHITE)]],
     anchor=MSO_ANCHOR.MIDDLE, space_after=4)
# Train/test distributions
text(s, Inches(0.6), Inches(5.15), Inches(12.0), Inches(0.5),
     [[("Train on a network-field distribution; test out-of-distribution on two held-out layouts:", 14, True, DARK)]])
pic_fit(s, os.path.join(IMG, "ic_distributions.png"), Inches(0.6), Inches(5.55), Inches(12.1), Inches(1.55))
footer(s, 5)

# ============================================================ 5b. MDP DIAGRAM
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

s = add_slide(); fill(s, WHITE)
header(s, "FORMULATION", "The control problem as an MDP")
text(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(0.55),
     [[("PhysiGym exposes the PhysiCell simulation as a Markov Decision Process — the SAC agent acts, the TME transitions, and a reward is returned each step.", 14, False, GREY, True)]],
     line_spacing=1.05)

def rounded(slide, l, t, w, h, fillc, linec=None, lw=Pt(1.5)):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fillc
    if linec is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = linec; sh.line.width = lw
    sh.shadow.inherit = False
    return sh

def settext(shape, lines, base=15):
    tf = shape.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, (txt, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = "Calibri"

# Agent box (top)
ag = rounded(s, Inches(4.5), Inches(2.35), Inches(4.3), Inches(1.25), ACCENT2)
settext(ag, [("AGENT", 20, True, WHITE), ("SAC policy  π(a | s)", 14, False, WHITE)])
# Environment box (bottom)
env = rounded(s, Inches(4.5), Inches(4.95), Inches(4.3), Inches(1.45), DARK)
settext(env, [("ENVIRONMENT", 20, True, WHITE),
              ("PhysiCell TME  ·  3 cell types · 5 substrates", 13, False, RGBColor(0xC8,0xCF,0xDA)),
              ("transition T,  reward R", 12, False, RGBColor(0xC8,0xCF,0xDA))])

def arrow(slide, x1, y1, x2, y2, color, width=Pt(3)):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color; cn.line.width = width
    cn.shadow.inherit = False
    le = cn.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    tail = le.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    le.append(tail)
    return cn

# action arrow: agent -> environment (right side, going down)
arrow(s, Inches(8.95), Inches(2.95), Inches(11.0), Inches(2.95), ACCENT)
arrow(s, Inches(11.0), Inches(2.95), Inches(11.0), Inches(5.65), ACCENT)
arrow(s, Inches(11.0), Inches(5.65), Inches(8.95), Inches(5.65), ACCENT)
text(s, Inches(9.1), Inches(3.55), Inches(3.0), Inches(1.5),
     [[("action  aₜ", 16, True, ACCENT)],
      [("(dose, x, y, r) ∈ [0,1]⁴", 12.5, False, RGBColor(0x33,0x3a,0x44))],
      [("localized drug_1", 12, False, GREY)]], align=PP_ALIGN.CENTER, space_after=3)

# state + reward arrow: environment -> agent (left side, going up)
arrow(s, Inches(4.35), Inches(5.65), Inches(2.3), Inches(5.65), ACCENT2)
arrow(s, Inches(2.3), Inches(5.65), Inches(2.3), Inches(2.95), ACCENT2)
arrow(s, Inches(2.3), Inches(2.95), Inches(4.35), Inches(2.95), ACCENT2)
text(s, Inches(1.2), Inches(3.55), Inches(3.0), Inches(1.6),
     [[("state  sₜ", 16, True, ACCENT2)],
      [("scalars / images", 12.5, False, RGBColor(0x33,0x3a,0x44))],
      [("reward  rₜ", 16, True, GREEN)],
      [("tumor suppression − dose", 11.5, False, GREY)]], align=PP_ALIGN.CENTER, space_after=2)

# objective strip
box(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.6), fill_color=LIGHT)
text(s, Inches(0.8), Inches(6.72), Inches(11.7), Inches(0.55),
     [[("Objective:   ", 14, True, DARK),
       ("maximise  𝔼[ Σ γᵗ rₜ ]  — eradicate the tumor while minimizing drug exposure.", 14, False, RGBColor(0x33,0x3a,0x44))]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 6)

# ============================================================ 6. OBSERVATION SPACES
s = add_slide(); fill(s, WHITE)
header(s, "STATE SPACE", "Six observation spaces compared")
text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.45),
     [[("From global scalars to full multi-channel spatial images  (N_c = 3 cell types, N_s = 5 substrates)", 14, False, GREY, True)]])
# scalar header
text(s, Inches(0.6), Inches(1.95), Inches(6), Inches(0.4), [[("Scalar observations", 16, True, ACCENT2)]])
text(s, Inches(7.0), Inches(1.95), Inches(6), Inches(0.4), [[("Image observations", 16, True, ACCENT)]])
scal = [
    ("S1  Normalised cell count", "Alive count per type vs. initial.   ℝ^Nc"),
    ("S2  Cells + substrate maxima", "S1 + max concentration per substrate.   ℝ^(Nc+Ns)"),
    ("S3  Per-type spatial statistics", "S1 + 6 spatial stats per cell type.   ℝ^7Nc"),
    ("S5  Global spatial statistics", "S3 + 6 global stats per substrate.   ℝ^(7Nc+6Ns)"),
]
y = Inches(2.45)
for h, b in scal:
    box(s, Inches(0.6), y, Inches(6.0), Inches(0.92), fill_color=LIGHT)
    text(s, Inches(0.8), y + Inches(0.06), Inches(5.7), Inches(0.85),
         [[(h, 13.5, True, DARK)], [(b, 12, False, RGBColor(0x33,0x3a,0x44))]], space_after=1)
    y += Inches(1.0)
img = [
    ("I1  Cell types only", "One channel per cell type.   ℝ^(Nc×64×64)  uint8"),
    ("I2  Cells + substrates", "I1 + all 5 substrate maps. Spatial feedback on drug & macrophage polarisation.   ℝ^((Nc+Ns)×64×64)"),
]
y = Inches(2.45)
for h, b in img:
    box(s, Inches(7.0), y, Inches(5.7), Inches(1.45), fill_color=RGBColor(0xFB,0xEC,0xEA))
    text(s, Inches(7.2), y + Inches(0.08), Inches(5.4), Inches(1.35),
         [[(h, 14, True, ACCENT)], [(b, 12.5, False, RGBColor(0x33,0x3a,0x44))]], space_after=2)
    y += Inches(1.6)
box(s, Inches(7.0), Inches(5.65), Inches(5.7), Inches(1.25), fill_color=DARK)
text(s, Inches(7.2), Inches(5.72), Inches(5.4), Inches(1.15),
     [[("Why substrate maps matter:", 12.5, True, ACCENT)],
      [("the indirect drug needs multi-step credit assignment. Only per-pixel substrate channels make polarisation state & drug distribution spatially legible.", 12, False, WHITE)]], space_after=2)
footer(s, 7)

# ============================================================ 7b. OBSERVATION EXAMPLE (I2)
s = add_slide(); fill(s, WHITE)
header(s, "STATE SPACE", "What the agent sees: the I2 observation")
text(s, Inches(0.6), Inches(1.35), Inches(5.3), Inches(0.55),
     [[("A single I2 observation = 8 stacked 64×64 channels.", 15, True, DARK)]], line_spacing=1.05)
obs = [
    ("3 cell channels", "tumor · T-cell · macrophage — where each population sits in the domain.", ACCENT2),
    ("5 substrate maps", "anti-tumoral & pro-tumoral factors, drug_1, tumor_molecule, cytokine.", GREEN),
    ("Drug_1 channel", "shows exactly where the agent injected — the disc visible in CH 5.", ACCENT),
    ("Closes the loop", "the CNN reads co-localisation (e.g. drug ∩ macrophages) directly from pixels — impossible from scalar summaries.", DARK),
]
y = Inches(2.0)
for h, b, c in obs:
    box(s, Inches(0.6), y, Inches(0.12), Inches(0.95), fill_color=c)
    text(s, Inches(0.85), y, Inches(5.0), Inches(1.0),
         [[(h, 14, True, c)], [(b, 12.5, False, RGBColor(0x33,0x3a,0x44))]], space_after=1)
    y += Inches(1.08)
pic_fit(s, os.path.join(IMG, "observation_40.png"), Inches(6.1), Inches(1.45), Inches(6.7), Inches(5.35))
text(s, Inches(6.1), Inches(6.75), Inches(6.7), Inches(0.4),
     [[("Example I2 observation at step 40 — 8 channels (3 cell types + 5 substrates), each 64×64.", 10.5, False, GREY, True)]], align=PP_ALIGN.CENTER)
footer(s, 8)

# ============================================================ 7. RESULTS TABLE
s = add_slide(); fill(s, WHITE)
header(s, "RESULTS", "Aggregated performance at end of training")
rows = [
    ("ID", "Mode", "Dim", "n", "Trainμ", "Trainσ", "Testμ", "Testσ"),
    ("I2", "img + substrates", "8×64²", "5", "+77.2", "8.5", "+56.5", "10.3"),
    ("I1", "img cells only", "3×64²", "4", "+50.3", "11.3", "+20.8", "18.2"),
    ("S3", "spatial cells", "21", "3", "+51.2", "25.8", "+12.4", "34.6"),
    ("S5", "spatial cells+subs", "39", "2", "+3.2", "7.3", "+0.5", "0.5"),
    ("S2", "scalar cells+subs", "8", "4", "−16.8", "38.3", "+12.3", "27.0"),
    ("S1", "scalar cells", "3", "4", "+5.9", "22.0", "+17.6", "22.5"),
]
from pptx.util import Cm
tbl_l, tbl_t = Inches(0.7), Inches(1.5)
tbl_w, tbl_h = Inches(8.2), Inches(3.9)
gtbl = s.shapes.add_table(len(rows), len(rows[0]), tbl_l, tbl_t, tbl_w, tbl_h).table
widths = [0.7, 2.3, 1.1, 0.6, 0.95, 0.95, 0.95, 0.95]
for i, w in enumerate(widths):
    gtbl.columns[i].width = Inches(w)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = gtbl.cell(ri, ci)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
        cell.margin_left = Pt(6); cell.margin_right = Pt(4)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if ci in (0,1) else PP_ALIGN.CENTER
        r = p.add_run(); r.text = val
        r.font.size = Pt(13)
        if ri == 0:
            r.font.bold = True; r.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = DARK
        else:
            highlight = (row[0] == "I2")
            r.font.bold = highlight
            r.font.color.rgb = ACCENT if highlight else RGBColor(0x2b,0x32,0x3c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFB,0xEC,0xEA) if highlight else (LIGHT if ri % 2 else WHITE)
# takeaways
box(s, Inches(9.2), Inches(1.5), Inches(3.5), Inches(3.9), fill_color=LIGHT)
text(s, Inches(9.4), Inches(1.6), Inches(3.15), Inches(3.7),
     [[("Key findings", 16, True, DARK)],
      [("• Substrate channels are decisive: I1→I2 gives the largest jump (+27 train / +36 test).", 12.5, False, RGBColor(0x33,0x3a,0x44))],
      [("• More scalar features doesn’t help — substrate scalars (S2) are even harmful.", 12.5, False, RGBColor(0x33,0x3a,0x44))],
      [("• Image modes show 3–4× lower seed variance (I2 σ=8.5 vs S2 σ=38.3).", 12.5, False, RGBColor(0x33,0x3a,0x44))]], space_after=8)
text(s, Inches(0.7), Inches(5.65), Inches(12), Inches(1.2),
     [[("Spatially explicit image observations including substrate maps learn faster and generalize better than scalar summaries.", 16, True, ACCENT)]], line_spacing=1.05)
footer(s, 9)

# ============================================================ 8. LEARNING CURVES
s = add_slide(); fill(s, WHITE)
header(s, "RESULTS", "Learning curves: training & generalization")
pic_fit(s, os.path.join(IMG, "train_return_mean50.png"), Inches(0.5), Inches(1.45), Inches(6.0), Inches(4.4))
pic_fit(s, os.path.join(IMG, "test_return_mean50.png"), Inches(6.85), Inches(1.45), Inches(6.0), Inches(4.4))
text(s, Inches(0.5), Inches(5.85), Inches(6.0), Inches(0.5),
     [[("Training return (EWMA-50). I2 (red) separates after ≈0.5×10⁵ steps and widens the gap; scalar modes plateau below +35.", 11, False, GREY, True)]], align=PP_ALIGN.CENTER)
text(s, Inches(6.85), Inches(5.85), Inches(6.0), Inches(0.5),
     [[("Test return on held-out layouts. I2 keeps climbing; I1 & scalar modes barely improve with high variance.", 11, False, GREY, True)]], align=PP_ALIGN.CENTER)
footer(s, 10)

# ============================================================ 9. GENERALIZATION TABLE + layouts
s = add_slide(); fill(s, WHITE)
header(s, "RESULTS", "Generalisation to held-out layouts")
rows = [
    ("Mode", "Test avg", "Rectangle", "Circular"),
    ("I2", "+56.5", "+44.5", "+97.9"),
    ("I1", "+20.8", "−0.6", "+5.0"),
    ("S3", "+12.4", "+112.7", "+16.3"),
    ("S5", "+0.5", "+85.1", "−37.6"),
    ("S2", "+12.3", "+40.2", "−47.7"),
    ("S1", "+17.6", "−7.7", "−20.5"),
]
gtbl = s.shapes.add_table(len(rows), 4, Inches(0.7), Inches(1.55), Inches(5.6), Inches(3.8)).table
for i, w in enumerate([1.3, 1.45, 1.45, 1.4]):
    gtbl.columns[i].width = Inches(w)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = gtbl.cell(ri, ci)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
        tf = cell.text_frame; p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = val; r.font.size = Pt(14)
        if ri == 0:
            r.font.bold = True; r.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = DARK
        else:
            hi = row[0] == "I2"
            r.font.bold = hi; r.font.color.rgb = ACCENT if hi else RGBColor(0x2b,0x32,0x3c)
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFB,0xEC,0xEA) if hi else (LIGHT if ri%2 else WHITE)
text(s, Inches(0.7), Inches(5.5), Inches(5.6), Inches(1.5),
     [[("I2 achieves the best average test return and the best circular return (+97.9).", 13, False, RGBColor(0x33,0x3a,0x44))],
      [("I1 fails on both OOD layouts — cell-image policies overfit without substrate feedback.", 13, False, RGBColor(0x33,0x3a,0x44))],
      [("Scalar modes generalize inconsistently (note the single-seed S3 rectangle outlier).", 13, False, RGBColor(0x33,0x3a,0x44))]], space_after=6)
pic_fit(s, os.path.join(IMG, "test_layout_comparison.png"), Inches(6.6), Inches(1.55), Inches(6.2), Inches(5.2))
footer(s, 11)

# ============================================================ 10. EPISODE COMPARISON
s = add_slide(); fill(s, WHITE)
header(s, "MECHANISM", "Step-by-step episode: why substrate maps win")
pic_fit(s, os.path.join(IMG, "fig_episode_comparison.png"), Inches(0.4), Inches(1.4), Inches(6.4), Inches(5.3))
text(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.5),
     [[("run_000143 — near-eradication episode", 16, True, DARK)]])
eps = [
    ("I2", "near-eradication: 4 tumor cells, +175 return, 35.4 dose.", ACCENT),
    ("S3", "partial control: 14 cells, +106 return, 31.1 dose.", ACCENT2),
    ("I1", "fails: 112 cells, −7.4 return, only 12.2 dose.", GREY),
]
y = Inches(2.1)
for h, b, c in eps:
    box(s, Inches(7.0), y, Inches(0.6), Inches(0.6), fill_color=c)
    text(s, Inches(7.0), y, Inches(0.6), Inches(0.6), [[(h, 14, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(7.75), y, Inches(5.0), Inches(0.7),
         [[(b, 13, False, RGBColor(0x33,0x3a,0x44))]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.85)
box(s, Inches(7.0), Inches(4.9), Inches(5.8), Inches(1.85), fill_color=LIGHT)
text(s, Inches(7.2), Inches(5.0), Inches(5.45), Inches(1.7),
     [[("Two failure modes:", 13.5, True, ACCENT)],
      [("I1 under-medicates — without substrate feedback it can’t tell if injections reached macrophages.", 12.5, False, RGBColor(0x33,0x3a,0x44))],
      [("S3 suppresses early but can’t finish — it lacks the substrate spatial feedback for the final push to eradication.", 12.5, False, RGBColor(0x33,0x3a,0x44))]], space_after=4)
footer(s, 12)

# ============================================================ 10b. VIDEOS
s = add_slide(); fill(s, WHITE)
header(s, "IN ACTION", "Learned policy vs. random policy")

def add_video(slide, path, poster, l, t, max_w, max_h):
    im = Image.open(poster); iw, ih = im.size; ar = iw/ih
    if max_w/max_h > ar:
        h = max_h; w = int(max_h*ar)
    else:
        w = max_w; h = int(max_w/ar)
    l2 = l + (max_w - w)//2
    mv = slide.shapes.add_movie(path, l2, t, w, h, poster_frame_image=poster,
                                mime_type="video/mp4")
    return mv

# left: merged best run (two best state spaces)
text(s, Inches(0.5), Inches(1.35), Inches(6.2), Inches(0.5),
     [[("Learned policy — two best state spaces (I2 & S3)", 16, True, ACCENT)]], align=PP_ALIGN.CENTER)
add_video(s, os.path.abspath(os.path.join(VID, "merged_run_000148.mp4")),
          os.path.join(IMG, "posters", "merged_poster.png"),
          Inches(0.5), Inches(1.95), Inches(6.2), Inches(3.9))
text(s, Inches(0.5), Inches(6.0), Inches(6.2), Inches(0.9),
     [[("run_000148 — side-by-side rollout of the two best observation modes. The agent targets drug injections to repolarise macrophages and clear a path for T-cells.", 12, False, RGBColor(0x33,0x3a,0x44))]],
     align=PP_ALIGN.CENTER, line_spacing=1.0)

# right: random policy
text(s, Inches(7.0), Inches(1.35), Inches(6.0), Inches(0.5),
     [[("Random policy — baseline", 16, True, GREY)]], align=PP_ALIGN.CENTER)
add_video(s, os.path.abspath(os.path.join(VID, "video.mp4")),
          os.path.join(IMG, "posters", "random_poster.png"),
          Inches(7.0), Inches(1.95), Inches(6.0), Inches(3.9))
text(s, Inches(7.0), Inches(6.0), Inches(6.0), Inches(0.9),
     [[("Untrained / random drug deployment — the tumor is not controlled, illustrating the difficulty of the indirect-control task.", 12, False, RGBColor(0x33,0x3a,0x44))]],
     align=PP_ALIGN.CENTER, line_spacing=1.0)
footer(s, 13)

# ============================================================ 11. TRAINING SETUP
s = add_slide(); fill(s, WHITE)
header(s, "TRAINING", "Asynchronous Soft Actor-Critic")
text(s, Inches(0.6), Inches(1.45), Inches(6.0), Inches(2.6),
     [[("Setup", 17, True, ACCENT2)],
      [("• SAC — sample-efficient continuous control; handles image obs via CNNs.", 14, False, RGBColor(0x33,0x3a,0x44))],
      [("• Asynchronous parallel framework: a GPU learner updates networks while 9 CPU env processes generate experience.", 14, False, RGBColor(0x33,0x3a,0x44))],
      [("• Scalar modes: MLP  FC(256)→Mish→FC(256).", 14, False, RGBColor(0x33,0x3a,0x44))],
      [("• Image modes: IMPALA-style CNN (3 conv → flatten → FC(256)×2).", 14, False, RGBColor(0x33,0x3a,0x44))]], space_after=8, line_spacing=1.05)
# hyperparam table
hp = [
    ("Parameter", "Value"),
    ("Total training steps", "5×10⁵"),
    ("Replay buffer size", "3×10⁵"),
    ("Batch size", "576"),
    ("Warm-up steps", "5,000"),
    ("Policy / critic LR", "3×10⁻⁴"),
    ("Discount γ", "0.99"),
    ("Target smoothing τ", "0.005"),
    ("Parallel envs", "9"),
    ("Update loops / step", "3"),
]
gtbl = s.shapes.add_table(len(hp), 2, Inches(7.0), Inches(1.45), Inches(5.6), Inches(5.0)).table
gtbl.columns[0].width = Inches(3.4); gtbl.columns[1].width = Inches(2.2)
for ri, row in enumerate(hp):
    for ci, val in enumerate(row):
        cell = gtbl.cell(ri, ci); cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Pt(1); cell.margin_bottom = Pt(1)
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = val; r.font.size = Pt(13)
        if ri == 0:
            r.font.bold=True; r.font.color.rgb=WHITE; cell.fill.solid(); cell.fill.fore_color.rgb=DARK
        else:
            r.font.color.rgb=RGBColor(0x2b,0x32,0x3c); cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if ri%2 else WHITE
text(s, Inches(7.0), Inches(6.55), Inches(5.6), Inches(0.4),
     [[("SAC hyperparameters (TME v2 experiments).", 10.5, False, GREY, True)]])
footer(s, 14)

# ============================================================ 12. CONCLUSION
s = add_slide(); fill(s, DARK)
box(s, 0, Inches(1.05), SW, Pt(4), fill_color=ACCENT)
text(s, Inches(0.8), Inches(0.35), SW-Inches(1.6), Inches(0.8),
     [[("Conclusion", 34, True, WHITE)]])
points = [
    ("PhysiGym bridges PhysiCell & Gymnasium", "direct RL policy learning in spatially explicit biological simulations, no recompilation, arbitrary obs/action spaces."),
    ("End-to-end case study works", "a SAC agent learns a spatial drug-delivery policy that drives the TME anti-tumoral and, at best, eradicates the tumor."),
    ("Substrate spatial images are the critical state", "+71 train / +39 test points over scalar baselines — a causal mechanism: they close the drug→repolarisation feedback loop."),
    ("A reproducible testbed", "for both the computational-biology and RL communities to study these questions principally."),
]
y = Inches(1.55)
for h, b in points:
    box(s, Inches(0.8), y+Inches(0.08), Inches(0.2), Inches(0.2), fill_color=ACCENT)
    text(s, Inches(1.2), y, Inches(11.4), Inches(1.05),
         [[(h, 18, True, WHITE), ("  —  "+b, 15, False, RGBColor(0xC8,0xCF,0xDA))]], line_spacing=1.05)
    y += Inches(1.15)
text(s, Inches(0.8), Inches(6.75), SW-Inches(1.6), Inches(0.5),
     [[("Open source:  https://anonymous.4open.science/r/PhysiGym-76A7", 14, True, ACCENT2)]])

prs.save(os.path.join(HERE, "PhysiGym_EWRL2026.pptx"))
print("Saved", os.path.join(HERE, "PhysiGym_EWRL2026.pptx"), "with", len(prs.slides._sldIdLst), "slides")
